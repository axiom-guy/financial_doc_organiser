#!/usr/bin/env python
# coding: utf-8

import os
import re
import shutil
import pandas as pd
import docx
from pptx import Presentation
import pdfplumber
import pytesseract

def read_text(path):
    try:
        with open(path,errors='ignore') as f:
            text=f.read(2000) # let 2000 be the max chars
        return text
    except Exception as e:
        print(f"Error reading text file {path}: {e}")
        return None

def read_docx(path):
    try:
        doc=docx.Document(path)
        text=[para.text for para in doc.paragraphs] 
        return '/n'.join(text)
    except Exception as e:
        print(f"Error reading docx file {path}: {e}")
        return None

def read_pdf(path):
    try:
        with pdfplumber.open(path) as f:
            return "\n".join(page.extract_text() or "" for page in f.pages)
    except Exception as e:
        print(f"Error reading pdf file {path}: {e}")
        return None

def read_spreadsheet(path):
    try:
        if path.lower().endswith('.csv'):
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
        text = df.to_string()
        return text
    except Exception as e:
        print(f"Error reading spreadsheet file {path}: {e}")
        return None

def read_ppt(path):
    try:
        prs=Presentation(path)
        f_text=[]
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f_text.append(shape.text)
        return '\n'.join(f_text)
    except Exception as e:
        print(f"Error reading ppt file {path}: {e}")
        return None

def read_image(path):
    try:
        from PIL import Image
        image=Image.open(path)
        text=pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error reading the image {path}: {e}")
    
def read_file(path):
    ext=os.path.splitext(path.lower())[1]
    if ext in ['.txt', '.md']:
        return read_text(path)
    elif ext in ['.docx', '.doc']:
        return read_docx(path)
    elif ext == '.pdf':
        return read_pdf(path)
    elif ext in ['.xls', '.xlsx', '.csv']:
        return read_spreadsheet(path)
    elif ext in ['.ppt', '.pptx']:
        return read_ppt(path)
    elif ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
        return read_image(path)
    else:
        return None

def collect_file_path(base_path):
    if os.path.isfile(base_path):
        if not os.path.basename(base_path).startswith('.'):
            return [base_path]
    else:
        file_paths=[]
        for root, _, files in os.walk(base_path):
            for file in files:
                if not file.startswith('.'): 
                    file_paths.append(os.path.join(root, file))
        return file_paths

